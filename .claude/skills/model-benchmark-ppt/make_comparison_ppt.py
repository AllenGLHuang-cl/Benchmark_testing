#!/usr/bin/env python3
"""Build a .pptx comparing image-generation models — two dataset layouts.

"features" layout (i2i / editing, one dir per feature):
    <feature>/src/<name>.<ext>          source images
    <feature>/meta_data/<name>.json     {"prompt": ["p0", "p1", ...]} (or shared meta_data/prompt.json)
    <feature>/results/<model_dir>/      model outputs, several naming conventions

"flat" layout (t2i / no source image):
    <model_dir>/                        flat dir per model, matched by key_regex

Usage:  python make_comparison_ppt.py config.json     (see example-config.json / example-config-flat.json)
"""
import io
import json
import math
import re
import sys
import datetime
from collections import Counter
from pathlib import Path

from PIL import Image, ImageFile
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ImageFile.LOAD_TRUNCATED_IMAGES = True

IMG_EXTS = (".jpg", ".jpeg", ".png", ".webp")

SLIDE_W = 13.333  # inches, 16:9
SLIDE_H = 7.5
MARGIN = 0.25
TITLE_H = 0.5
PROMPT_H = 0.45
LABEL_H = 0.26
CELL_PAD = 0.06

GREY_FILL = RGBColor(0xEB, 0xEB, 0xEB)
GREY_TEXT = RGBColor(0x99, 0x99, 0x99)
DARK_TEXT = RGBColor(0x33, 0x33, 0x33)
PERF_ORANGE = RGBColor(0xF8, 0xCB, 0xAD)


# ---------------------------------------------------------------- resolver

def resolve_output(model_dir: Path, src: str, p: int, seed_prefer: list):
    """Find the output image for (src stem, prompt index) in a model dir.

    Tries the known naming conventions in order; returns Path or None.
      1. {src}_prompt{p}_seed{S}.<ext> or {src}_prompt{p}_out{i}_seed{S}.<ext>
         (pick by seed_prefer, then lexicographic)
      2. {src}_prompt{p}.<ext>
      3. {src}_p{p}.<ext>
      4. {src}_{p}_{n}.<ext>             (pick smallest n)
      5. {src}/p{p}_{n}.<ext>            (per-src subdirectory, smallest n)
    """
    if not model_dir.is_dir():
        return None
    esc = re.escape(src)
    ext = r"\.(?:jpg|jpeg|png|webp)$"
    names = [f.name for f in model_dir.iterdir() if f.is_file()]

    def pick(regex, key=None):
        rx = re.compile(regex, re.IGNORECASE)
        hits = [(rx.match(n), n) for n in names]
        hits = [(m, n) for m, n in hits if m]
        if not hits:
            return None
        hits.sort(key=key if key else lambda h: h[1])
        return model_dir / hits[0][1]

    def seed_rank(hit):
        m, name = hit
        tag = "seed" + m.group(1)
        return (seed_prefer.index(tag) if tag in seed_prefer else len(seed_prefer), name)

    found = pick(rf"^{esc}_prompt{p}(?:_out\d+)?_seed(\w+){ext}", key=seed_rank)
    if found:
        return found
    found = pick(rf"^{esc}_prompt{p}{ext}")
    if found:
        return found
    found = pick(rf"^{esc}_p{p}{ext}")
    if found:
        return found
    found = pick(rf"^{esc}_{p}_(\d+){ext}", key=lambda h: int(h[0].group(1)))
    if found:
        return found
    sub = model_dir / src
    if sub.is_dir():
        rx = re.compile(rf"^p{p}_(\d+){ext}", re.IGNORECASE)
        hits = sorted(
            ((rx.match(f.name), f) for f in sub.iterdir() if f.is_file() and rx.match(f.name)),
            key=lambda h: int(h[0].group(1)),
        )
        if hits:
            return hits[0][1]
    return None


# ---------------------------------------------------------------- groups

def enumerate_groups(feature_dir: Path, picks):
    """Yield (src_stem, src_path, prompt_idx, prompt_text); warn+skip broken entries.

    Prompts come from meta_data/{stem}.json, falling back to the shared
    meta_data/prompt.json (same two-tier lookup as websocket_edit.py's
    load_prompts_for_image) — most features use the shared file, so stems
    must be enumerated from src/, not from meta_data/*.json."""
    meta_dir = feature_dir / "meta_data"
    src_dir = feature_dir / "src"
    shared_meta = meta_dir / "prompt.json"
    for src_path in sorted(f for f in src_dir.iterdir() if f.is_file() and f.suffix.lower() in IMG_EXTS):
        stem = src_path.stem
        if picks is not None and stem not in picks:
            continue
        meta = meta_dir / f"{stem}.json"
        if not meta.is_file():
            meta = shared_meta
            if not meta.is_file():
                print(f"  WARN [{feature_dir.name}] no meta_data for '{stem}' — skipped")
                continue
        try:
            prompts = json.loads(meta.read_text())["prompt"]
        except (json.JSONDecodeError, KeyError, TypeError) as e:
            print(f"  WARN [{feature_dir.name}] bad {meta.name}: {e} — skipped")
            continue
        if isinstance(prompts, str):
            prompts = [prompts]
        for p_idx, p_text in enumerate(prompts):
            yield stem, src_path, p_idx, p_text


# ---------------------------------------------------------------- flat (t2i) layout

def index_model(d: Path, key_regex, prefer_regex):
    """dir -> {key: best_path}. key = group(1) of key_regex applied to basename.
    If several files share a key, prefer one matching prefer_regex, else first sorted."""
    if not d.is_dir():
        return {}
    pat = re.compile(key_regex)
    pref = re.compile(prefer_regex) if prefer_regex else None
    buckets = {}
    for p in sorted(d.iterdir()):
        if not p.is_file():
            continue
        m = pat.search(p.name)
        if m:
            buckets.setdefault(m.group(1), []).append(p)
    out = {}
    for k, paths in buckets.items():
        chosen = paths[0]
        if pref:
            for p in paths:
                if pref.search(p.name):
                    chosen = p
                    break
        out[k] = chosen
    return out


def sort_keys(keys):
    try:
        return sorted(keys, key=lambda x: int(x))
    except ValueError:
        return sorted(keys)


def load_captions(cfg):
    c = cfg.get("captions")
    if not c:
        return {}
    data = json.loads(Path(c["file"]).read_text(encoding="utf-8"))
    if c.get("json_path"):
        data = data[c["json_path"]]
    if c.get("key_is_index"):  # data is a list; caption[str(i)] = data[i]
        return {str(i): str(t) for i, t in enumerate(data)}
    return {str(k): str(v) for k, v in data.items()}


def build_flat_groups(cfg, models):
    """t2i-style layout: models[].dir is a flat directory, matched by key_regex.
    Only keys present in every model are kept (same intersect-don't-guess philosophy
    as the features layout dropping incomplete groups)."""
    key_regex = cfg["key_regex"]
    prefer = cfg.get("prefer_regex")

    indexed = [index_model(Path(m["dir"]), key_regex, prefer) for m in models]
    common = set(indexed[0])
    for ix in indexed[1:]:
        common &= set(ix)
    common = sort_keys(common)

    if cfg.get("picks"):
        picks = [str(k) for k in cfg["picks"]]
        missing = [k for k in picks if k not in common]
        if missing:
            print(f"WARNING picks not common to all models (skipped): {missing}")
        picks = [k for k in picks if k in common]
    else:
        n = cfg.get("auto_pick", min(30, len(common)))
        if n >= len(common):
            picks = common
        else:  # evenly spaced sample across the common set
            step = len(common) / n
            picks = [common[int(i * step)] for i in range(n)]
    print(f"{len(common)} common keys; rendering {len(picks)} groups across {len(models)} models")

    captions = load_captions(cfg)
    return [{"key": k, "caption": captions.get(k, ""), "images": [ix[k] for ix in indexed]} for k in picks]


# ---------------------------------------------------------------- images

def load_thumb(path: Path, thumb_px: int):
    """Return (BytesIO jpeg, w, h) downscaled to thumb_px long side, or None if unreadable."""
    try:
        with Image.open(path) as im:
            im = im.convert("RGB")
            im.thumbnail((thumb_px, thumb_px))
            buf = io.BytesIO()
            im.save(buf, "JPEG", quality=85)
            buf.seek(0)
            return buf, im.width, im.height
    except OSError as e:
        print(f"  WARN unreadable image {path}: {e}")
        return None


# ---------------------------------------------------------------- layout

def best_grid(n_cells, region_w, region_h, reserved_cols=0):
    """cols, rows maximizing the square image side that fits every cell.
    reserved_cols: extra same-width columns set aside (e.g. a dedicated Source
    column on the left) — they consume width but hold none of the n_cells."""
    best = (1, n_cells, 0.0)
    for cols in range(1, n_cells + 1):
        rows = math.ceil(n_cells / cols)
        cw = region_w / (cols + reserved_cols)
        ch = region_h / rows
        side = min(cw - 2 * CELL_PAD, ch - LABEL_H - 2 * CELL_PAD)
        if side > best[2]:
            best = (cols, rows, side)
    return best[0], best[1]


def add_text(slide, x, y, w, h, text, size, bold=False, italic=False,
             color=DARK_TEXT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def place_cell(slide, x, y, cw, ch, label, thumb):
    """One grid cell: image (or grey 'missing' box) + label underneath."""
    img_h = ch - LABEL_H - 2 * CELL_PAD
    img_w = cw - 2 * CELL_PAD
    if thumb is not None:
        buf, w_px, h_px = thumb
        scale = min(img_w / w_px, img_h / h_px)
        w_in, h_in = w_px * scale, h_px * scale
        slide.shapes.add_picture(
            buf,
            Inches(x + (cw - w_in) / 2),
            Inches(y + CELL_PAD + (img_h - h_in) / 2),
            Inches(w_in), Inches(h_in),
        )
    else:
        box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x + CELL_PAD), Inches(y + CELL_PAD), Inches(img_w), Inches(img_h),
        )
        box.fill.solid()
        box.fill.fore_color.rgb = GREY_FILL
        box.line.color.rgb = GREY_TEXT
        tf = box.text_frame
        tf.text = "missing"
        tf.paragraphs[0].runs[0].font.color.rgb = GREY_TEXT
        tf.paragraphs[0].runs[0].font.size = Pt(12)
    add_text(slide, x, y + ch - LABEL_H, cw, LABEL_H, label, 10,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def render_group_slide(prs, cols, cw, ch, region_y, region_w, title, prompt_text, cell_specs, thumb_px,
                       source_left=False):
    """One slide: title + prompt + grid of cell_specs=[(label, path_or_None), ...].
    Shared by both dataset layouts. Returns the labels whose image was missing/unreadable.
    source_left: cell_specs[0] (the Source) gets its own left column, vertically
    centered; the remaining cells grid into the columns to its right."""
    slide = blank_slide(prs)
    add_text(slide, MARGIN, MARGIN, region_w, TITLE_H, title, 18, bold=True)
    shown = prompt_text if len(prompt_text) <= 220 else prompt_text[:217] + "…"
    box = add_text(slide, MARGIN, MARGIN + TITLE_H, region_w, PROMPT_H,
                   "Prompt: ", 12, bold=True)
    run = box.text_frame.paragraphs[0].add_run()
    run.text = f"“{shown}”"
    run.font.size = Pt(12)
    run.font.italic = True
    run.font.color.rgb = GREY_TEXT

    notes = [f"prompt: {prompt_text}"]
    missing_labels = []
    for i, (label, path) in enumerate(cell_specs):
        thumb = load_thumb(path, thumb_px) if path else None
        if thumb is None:
            missing_labels.append(label)
            notes.append(f"{label}: MISSING")
        else:
            notes.append(f"{label}: {path}")
        if source_left and i == 0:
            region_h = SLIDE_H - region_y - MARGIN
            place_cell(slide, MARGIN, region_y + (region_h - ch) / 2, cw, ch, label, thumb)
        elif source_left:
            r, c = divmod(i - 1, cols)
            place_cell(slide, MARGIN + (c + 1) * cw, region_y + r * ch, cw, ch, label, thumb)
        else:
            r, c = divmod(i, cols)
            place_cell(slide, MARGIN + c * cw, region_y + r * ch, cw, ch, label, thumb)
    slide.notes_slide.notes_text_frame.text = "\n".join(notes)
    return missing_labels


# ---------------------------------------------------------------- perf table

# performance.jsonl record keys (written by websocket_edit.py / websocket_tti.py,
# see run_batch/perf_utils.py — record schema kept in sync by hand across the two,
# since this skill stays self-contained rather than importing run_batch/) ->
# the aggregate PPT column each one feeds.
_PERF_AUTO_COLUMNS = {
    "inference steps": ("steps", "mode"),
    "inference time (s)": ("elapsed_sec", "avg"),
    "memory usage (GiB)": ("vram_peak_gib", "peak"),
    "weight_type": ("weight_type", "mode"),
    "device": ("device", "mode"),
    "n_samples": (None, "count"),
}


def aggregate_perf_jsonl(path: Path):
    """Read a performance.jsonl and reduce it to one row of PPT columns (see
    _PERF_AUTO_COLUMNS). Unknown/absent fields are simply omitted — never guessed."""
    records = []
    for line in path.read_text().splitlines():
        line = line.strip()
        if line:
            records.append(json.loads(line))
    if not records:
        return {}

    def values(key):
        return [r[key] for r in records if r.get(key) is not None]

    row = {}
    for col, (key, agg) in _PERF_AUTO_COLUMNS.items():
        if agg == "count":
            row[col] = len(records)
            continue
        vals = values(key)
        if not vals:
            continue
        if agg == "mode":
            row[col] = Counter(vals).most_common(1)[0][0]
        elif agg == "avg":
            row[col] = round(sum(vals) / len(vals), 1)
        elif agg == "peak":
            row[col] = round(max(vals), 1)
    w, h = values("output_width"), values("output_height")
    if w and h:
        row["output_size"] = f"{Counter(w).most_common(1)[0][0]}x{Counter(h).most_common(1)[0][0]}"
    return row


def build_perf_rows(perf_cfg, models, root, layout):
    """Merge auto-computed performance.jsonl aggregates (if perf.auto_from is set)
    with any hand-typed perf.rows — manual values win so specific cells can be
    corrected/added (e.g. 'Model size' which can't be measured)."""
    auto_feature = perf_cfg.get("auto_from")
    manual_rows = perf_cfg.get("rows", {})
    rows = {}
    for m in models:
        label = m["label"]
        row = {}
        if auto_feature:
            log_path = (Path(m["dir"]) / "performance.jsonl" if layout == "flat"
                        else root / auto_feature / "results" / m["dir"] / "performance.jsonl")
            if log_path.is_file():
                row.update(aggregate_perf_jsonl(log_path))
            else:
                print(f"WARN perf: no performance.jsonl at {log_path} for '{label}' — leaving auto columns blank")
        row.update(manual_rows.get(label) or manual_rows.get(m["dir"]) or {})
        rows[label] = row
    return rows


def add_perf_slide(prs, perf, models, rows_cfg):
    """Quantitative table: rows = models, columns from config. Unknown cells stay blank."""
    from pptx.oxml.ns import qn

    slide = blank_slide(prs)
    add_text(slide, MARGIN, 0.45, SLIDE_W - 2 * MARGIN, 0.6,
             perf.get("note", "Performance"), 20, color=DARK_TEXT, align=PP_ALIGN.CENTER)
    cols = perf["columns"]
    n_rows, n_cols = len(models) + 1, len(cols) + 1
    first_w, col_w = 2.6, 2.0
    tbl_w = min(first_w + col_w * len(cols), SLIDE_W - 2 * MARGIN)
    tbl_h = 0.55 * n_rows
    frame = slide.shapes.add_table(
        n_rows, n_cols,
        Inches((SLIDE_W - tbl_w) / 2), Inches(max(1.6, (SLIDE_H - tbl_h) / 2)),
        Inches(tbl_w), Inches(tbl_h),
    )
    table = frame.table
    # "No Style, Table Grid" — plain black grid instead of the default blue banding
    table._tbl.tblPr.find(qn('a:tableStyleId')).text = \
        '{5940675A-B579-460E-94D1-54222C63F5DA}'
    table.columns[0].width = Inches(first_w)
    for j in range(len(cols)):
        table.columns[j + 1].width = Inches((tbl_w - first_w) / len(cols))

    def set_cell(r, c, text, fill=None):
        cell = table.cell(r, c)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        if fill is not None:
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill
        para = cell.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = text
        run.font.size = Pt(14)
        run.font.bold = False
        run.font.color.rgb = DARK_TEXT

    set_cell(0, 0, "", fill=PERF_ORANGE)
    for j, col in enumerate(cols):
        set_cell(0, j + 1, col)
    for i, m in enumerate(models):
        vals = rows_cfg.get(m["label"], {})
        set_cell(i + 1, 0, m["label"])
        for j, col in enumerate(cols):
            v = vals.get(col)
            set_cell(i + 1, j + 1, "" if v is None else str(v))


# ---------------------------------------------------------------- deck

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def build(cfg_path: Path):
    cfg = json.loads(cfg_path.read_text())
    models = cfg["models"]
    thumb_px = cfg.get("thumb_px", 768)
    layout = cfg.get("dataset_layout", "features")
    root = Path(cfg["dataset_root"]) if "dataset_root" in cfg else None

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    # cover
    slide = blank_slide(prs)
    add_text(slide, 1.0, 2.4, SLIDE_W - 2.0, 1.0,
             cfg.get("title", "Model Comparison"), 36, bold=True)
    sub = (datetime.date.today().isoformat() + "  ·  " + (root.name if root else layout) + "\n"
           + "Models: " + ", ".join(m["label"] for m in models))
    add_text(slide, 1.0, 3.5, SLIDE_W - 2.0, 1.5, sub, 16, color=GREY_TEXT)

    if cfg.get("perf"):
        rows_cfg = build_perf_rows(cfg["perf"], models, root, layout)
        add_perf_slide(prs, cfg["perf"], models, rows_cfg)

    region_w = SLIDE_W - 2 * MARGIN
    region_y = MARGIN + TITLE_H + PROMPT_H
    region_h = SLIDE_H - region_y - MARGIN
    n_groups = 0
    missing = {m["label"]: 0 for m in models}

    if layout == "flat":
        cols, rows = best_grid(len(models), region_w, region_h)
        cw, ch = region_w / cols, region_h / rows
        for g in build_flat_groups(cfg, models):
            n_groups += 1
            cell_specs = list(zip((m["label"] for m in models), g["images"]))
            for label in render_group_slide(prs, cols, cw, ch, region_y, region_w,
                                             g["key"], g["caption"], cell_specs, thumb_px):
                missing[label] += 1
    else:
        seed_prefer = cfg.get("seed_prefer", ["seed2025"])
        picks_cfg = cfg.get("picks", {})

        features = cfg["features"]
        if features == "all":
            features = sorted(
                d.name for d in root.iterdir()
                if d.is_dir() and not d.name.startswith("_")
                and (d / "meta_data").is_dir() and (d / "results").is_dir()
            )
        kept = []
        for feat in features:
            fdir = root / feat
            if not (fdir / "meta_data").is_dir() or not (fdir / "results").is_dir():
                print(f"WARN feature '{feat}' has no meta_data/results — skipped")
                continue
            kept.append(feat)
            for m in models:
                if not (fdir / "results" / m["dir"]).is_dir():
                    print(f"WARN [{feat}] no results dir for model '{m['dir']}' — whole column missing")

        cols, rows = best_grid(len(models), region_w, region_h, reserved_cols=1)
        cw, ch = region_w / (cols + 1), region_h / rows

        for feat in kept:
            fdir = root / feat
            groups = list(enumerate_groups(fdir, picks_cfg.get(feat)))
            if not groups:
                print(f"WARN feature '{feat}' produced no comparison groups")
                continue

            slide = blank_slide(prs)  # section divider
            add_text(slide, 1.0, 3.0, SLIDE_W - 2.0, 1.2, feat, 40, bold=True,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            add_text(slide, 1.0, 4.2, SLIDE_W - 2.0, 0.5, f"{len(groups)} comparisons",
                     16, color=GREY_TEXT, align=PP_ALIGN.CENTER)

            for stem, src_path, p_idx, p_text in groups:
                n_groups += 1
                cell_specs = [("Source", src_path)]
                for m in models:
                    cell_specs.append((m["label"], resolve_output(fdir / "results" / m["dir"], stem, p_idx, seed_prefer)))
                title = f"{feat} / {stem}  [p{p_idx}]"
                for label in render_group_slide(prs, cols, cw, ch, region_y, region_w,
                                                 title, p_text, cell_specs, thumb_px,
                                                 source_left=True):
                    if label != "Source":
                        missing[label] += 1

    out = Path(cfg["out"])
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)

    print(f"\nDone: {out}  ({out.stat().st_size / 1e6:.1f} MB)")
    print(f"  layout: {layout}   comparison groups: {n_groups}   grid: {cols}x{rows}")
    for label, n in missing.items():
        print(f"  missing cells — {label}: {n}/{n_groups}")


if __name__ == "__main__":
    if len(sys.argv) != 2:
        sys.exit(__doc__)
    build(Path(sys.argv[1]))
